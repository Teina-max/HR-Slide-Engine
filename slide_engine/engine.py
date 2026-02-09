"""Core engine functions for HR Slide Engine."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from . import design as D


def create_presentation():
    """Create a new 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = D.SLIDE_WIDTH
    prs.slide_height = D.SLIDE_HEIGHT
    return prs


def save_presentation(prs, filename):
    """Save presentation to file. Appends .pptx if missing."""
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    prs.save(filename)
    return filename


def _add_blank_slide(prs):
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def _set_slide_background(slide, color):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text,
                 font_size=D.BODY_SIZE, font_color=D.DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=D.FONT_FAMILY, anchor=MSO_ANCHOR.TOP):
    """Add a textbox with formatted text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Set vertical alignment
    txBox.text_frame._txBody.bodyPr.set("anchor", {
        MSO_ANCHOR.TOP: "t",
        MSO_ANCHOR.MIDDLE: "ctr",
        MSO_ANCHOR.BOTTOM: "b",
    }.get(anchor, "t"))

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return txBox


def _add_multiline_textbox(slide, left, top, width, height, lines,
                           font_size=D.BODY_SIZE, font_color=D.DARK_TEXT,
                           bold=False, alignment=PP_ALIGN.LEFT,
                           font_name=D.FONT_FAMILY, line_spacing=None,
                           bullet_color=None, bullet_char=None):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet_char:
            run_bullet = p.add_run()
            run_bullet.text = f"{bullet_char} "
            run_bullet.font.size = font_size
            run_bullet.font.bold = bold
            run_bullet.font.name = font_name
            run_bullet.font.color.rgb = bullet_color or font_color

            run_text = p.add_run()
            run_text.text = line
            run_text.font.size = font_size
            run_text.font.bold = False
            run_text.font.name = font_name
            run_text.font.color.rgb = font_color
        else:
            p.text = line
            p.font.size = font_size
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.font.name = font_name

        p.alignment = alignment
        if line_spacing:
            p.space_after = line_spacing

    return txBox


def _add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def _add_rectangle(slide, left, top, width, height, fill_color):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # No border
    return shape


def _add_line(slide, left, top, width, height, color, line_width=Pt(2)):
    """Add a line shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_rounded_rectangle(slide, left, top, width, height, fill_color,
                           border_color=None, text="", font_size=D.BODY_SIZE,
                           font_color=D.WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    """Add a rounded rectangle with optional text inside."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = D.FONT_FAMILY
        tf.paragraphs[0].alignment = alignment
        shape.text_frame._txBody.bodyPr.set("anchor", "ctr")
    return shape


def _add_chevron(slide, left, top, width, height, fill_color, text="",
                 font_size=D.SMALL_SIZE, font_color=D.WHITE):
    """Add a chevron (pentagon/arrow) shape with text."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = D.FONT_FAMILY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame._txBody.bodyPr.set("anchor", "ctr")
    return shape


def _add_oval(slide, left, top, width, height, fill_color, text="",
              font_size=D.BODY_SIZE, font_color=D.WHITE, bold=True):
    """Add an oval/circle shape with text."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = D.FONT_FAMILY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame._txBody.bodyPr.set("anchor", "ctr")
    return shape


def _add_triangle(slide, left, top, width, height, fill_color):
    """Add an isoceles triangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_chart_bar(slide, left, top, width, height, categories, values,
                   chart_title=""):
    """Add a bar chart to the slide."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Style the bars with orange
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = D.ORANGE

    # Style axes
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)
    category_axis.tick_labels.font.name = D.FONT_FAMILY
    category_axis.tick_labels.font.color.rgb = D.GRAY

    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(11)
    value_axis.tick_labels.font.name = D.FONT_FAMILY
    value_axis.tick_labels.font.color.rgb = D.GRAY
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = D.LIGHT_GRAY

    return chart_frame


def _add_chart_pie(slide, left, top, width, height, categories, values):
    """Add a pie chart to the slide."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    # Color each slice
    plot = chart.plots[0]
    colors = D.PROCESS_COLORS
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]

    # Data labels with percentages
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_category_name = True
    data_labels.show_value = False
    data_labels.font.size = Pt(11)
    data_labels.font.name = D.FONT_FAMILY
    data_labels.font.color.rgb = D.DARK_TEXT

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = D.FONT_FAMILY
    chart.legend.include_in_layout = False

    return chart_frame
