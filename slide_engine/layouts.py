"""Layout functions for HR Slide Engine — 18 professional slide types."""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import design as D
from .engine import (
    _add_blank_slide,
    _set_slide_background,
    _add_textbox,
    _add_multiline_textbox,
    _add_speaker_notes,
    _add_rectangle,
    _add_line,
    _add_rounded_rectangle,
    _add_chevron,
    _add_oval,
    _add_triangle,
    _add_chart_bar,
    _add_chart_pie,
)


def add_title_slide(prs, title, subtitle="", notes=""):
    """Slide 1 — Title: navy background, white centered text."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.NAVY)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=Inches(2.2),
        width=D.CONTENT_WIDTH, height=Inches(1.5),
        text=title,
        font_size=Pt(36), font_color=D.WHITE,
        bold=True, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.BOTTOM,
    )

    # Orange accent line
    line_width = Inches(3)
    line_left = (D.SLIDE_WIDTH - line_width) // 2
    _add_line(slide, line_left, Inches(3.8), line_width, Pt(3), D.ORANGE)

    # Subtitle
    if subtitle:
        _add_textbox(
            slide,
            left=D.MARGIN_LEFT, top=Inches(4.1),
            width=D.CONTENT_WIDTH, height=Inches(1.0),
            text=subtitle,
            font_size=D.SUBTITLE_SIZE, font_color=D.LIGHT_GRAY,
            bold=False, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.TOP,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_agenda_slide(prs, items, title="Agenda", notes=""):
    """Slide 2 — Agenda: numbered list with orange numbers."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    # Underline
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    # Items
    y_start = Inches(2.0)
    item_height = Inches(0.6)
    for i, item in enumerate(items, 1):
        y = y_start + (i - 1) * item_height

        # Orange number
        _add_textbox(
            slide,
            left=D.MARGIN_LEFT, top=y,
            width=Inches(0.6), height=item_height,
            text=f"{i:02d}",
            font_size=Pt(22), font_color=D.ORANGE,
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # Item text
        _add_textbox(
            slide,
            left=D.MARGIN_LEFT + Inches(0.7), top=y,
            width=D.CONTENT_WIDTH - Inches(0.7), height=item_height,
            text=item,
            font_size=D.BODY_SIZE, font_color=D.DARK_TEXT,
            bold=False, alignment=PP_ALIGN.LEFT,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_section_slide(prs, title, subtitle="", notes=""):
    """Slide 3 — Section divider: navy bar on the left, large title."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Navy vertical bar
    _add_rectangle(
        slide,
        left=Inches(0.4), top=Inches(1.5),
        width=D.SECTION_BAR_WIDTH, height=Inches(4.5),
        fill_color=D.NAVY,
    )

    # Title
    _add_textbox(
        slide,
        left=Inches(1.0), top=Inches(2.5),
        width=Inches(10.5), height=Inches(1.5),
        text=title,
        font_size=Pt(32), font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.BOTTOM,
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide,
            left=Inches(1.0), top=Inches(4.2),
            width=Inches(10.5), height=Inches(0.8),
            text=subtitle,
            font_size=D.SUBTITLE_SIZE, font_color=D.GRAY,
            bold=False, alignment=PP_ALIGN.LEFT,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_bullets_slide(prs, title, bullets, notes=""):
    """Slide 4 — Bullet points: orange bullets, gray text."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    # Underline
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    # Bullet items
    _add_multiline_textbox(
        slide,
        left=D.MARGIN_LEFT + Inches(0.3), top=Inches(2.0),
        width=D.CONTENT_WIDTH - Inches(0.3), height=Inches(4.5),
        lines=bullets,
        font_size=D.BODY_SIZE, font_color=D.GRAY,
        bullet_char=D.BULLET_CHAR, bullet_color=D.ORANGE,
        line_spacing=D.PARAGRAPH_SPACING,
    )

    _add_speaker_notes(slide, notes)
    return slide


def add_two_columns_slide(prs, title, left_title, left_items,
                          right_title, right_items, notes=""):
    """Slide 5 — Two columns: separated by a thin gray line."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    col_width = (D.CONTENT_WIDTH - D.COLUMN_GAP) / 2
    left_x = D.MARGIN_LEFT
    right_x = D.MARGIN_LEFT + col_width + D.COLUMN_GAP

    # Vertical separator
    sep_x = D.MARGIN_LEFT + col_width + (D.COLUMN_GAP // 2)
    _add_line(slide, sep_x, Inches(2.0), Pt(1), Inches(4.5), D.LIGHT_GRAY)

    # Left column title
    _add_textbox(
        slide,
        left=left_x, top=Inches(2.0),
        width=col_width, height=Inches(0.6),
        text=left_title,
        font_size=Pt(22), font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    # Left column items
    _add_multiline_textbox(
        slide,
        left=left_x + Inches(0.2), top=Inches(2.7),
        width=col_width - Inches(0.2), height=Inches(3.8),
        lines=left_items,
        font_size=Pt(18), font_color=D.GRAY,
        bullet_char=D.BULLET_CHAR, bullet_color=D.ORANGE,
        line_spacing=D.LINE_SPACING,
    )

    # Right column title
    _add_textbox(
        slide,
        left=right_x, top=Inches(2.0),
        width=col_width, height=Inches(0.6),
        text=right_title,
        font_size=Pt(22), font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )

    # Right column items
    _add_multiline_textbox(
        slide,
        left=right_x + Inches(0.2), top=Inches(2.7),
        width=col_width - Inches(0.2), height=Inches(3.8),
        lines=right_items,
        font_size=Pt(18), font_color=D.GRAY,
        bullet_char=D.BULLET_CHAR, bullet_color=D.ORANGE,
        line_spacing=D.LINE_SPACING,
    )

    _add_speaker_notes(slide, notes)
    return slide


def add_key_stat_slide(prs, stat, description, notes=""):
    """Slide 6 — Key statistic: large orange number centered."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Big stat
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=Inches(1.8),
        width=D.CONTENT_WIDTH, height=Inches(2.5),
        text=stat,
        font_size=D.STAT_SIZE, font_color=D.ORANGE,
        bold=True, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.BOTTOM,
    )

    # Description
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=Inches(4.5),
        width=D.CONTENT_WIDTH, height=Inches(1.5),
        text=description,
        font_size=D.BODY_SIZE, font_color=D.GRAY,
        bold=False, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
    )

    _add_speaker_notes(slide, notes)
    return slide


def add_quote_slide(prs, quote, author="", notes=""):
    """Slide 7 — Quote: light gray background, decorative quotation mark."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.LIGHT_GRAY)

    # Decorative large quote mark
    _add_textbox(
        slide,
        left=Inches(1.0), top=Inches(1.0),
        width=Inches(2.0), height=Inches(2.0),
        text=D.QUOTE_CHAR,
        font_size=Pt(120), font_color=D.ORANGE,
        bold=False, alignment=PP_ALIGN.LEFT,
    )

    # Quote text
    _add_textbox(
        slide,
        left=Inches(2.0), top=Inches(2.5),
        width=Inches(9.0), height=Inches(2.5),
        text=quote,
        font_size=D.QUOTE_SIZE, font_color=D.NAVY,
        bold=False, alignment=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Author
    if author:
        _add_textbox(
            slide,
            left=Inches(2.0), top=Inches(5.3),
            width=Inches(9.0), height=Inches(0.6),
            text=f"— {author}",
            font_size=D.SUBTITLE_SIZE, font_color=D.GRAY,
            bold=False, alignment=PP_ALIGN.LEFT,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_conclusion_slide(prs, title, points, notes=""):
    """Slide 8 — Conclusion: navy banner at top, checkmarks orange."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Navy banner
    _add_rectangle(
        slide,
        left=Inches(0), top=Inches(0),
        width=D.SLIDE_WIDTH, height=Inches(1.8),
        fill_color=D.NAVY,
    )

    # Title on banner
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=Inches(0.4),
        width=D.CONTENT_WIDTH, height=Inches(1.0),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.WHITE,
        bold=True, alignment=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Conclusion points with checkmarks
    _add_multiline_textbox(
        slide,
        left=D.MARGIN_LEFT + Inches(0.3), top=Inches(2.3),
        width=D.CONTENT_WIDTH - Inches(0.3), height=Inches(4.5),
        lines=points,
        font_size=D.BODY_SIZE, font_color=D.DARK_TEXT,
        bullet_char=D.CHECKMARK_CHAR, bullet_color=D.ORANGE,
        line_spacing=D.PARAGRAPH_SPACING,
    )

    _add_speaker_notes(slide, notes)
    return slide


# ===================================================================
# VISUAL LAYOUTS — Diagrams, charts, and graphic slides
# ===================================================================


def add_process_flow_slide(prs, title, steps, notes=""):
    """Slide 9 — Process flow: connected chevron arrows, colored steps."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    n = len(steps)
    total_width = D.CONTENT_WIDTH
    gap = Inches(0.05)
    chevron_w = (total_width - gap * (n - 1)) / n
    chevron_h = Inches(1.2)
    y_chevron = Inches(2.5)
    colors = D.PROCESS_COLORS

    for i, step in enumerate(steps):
        x = D.MARGIN_LEFT + i * (chevron_w + gap)
        color = colors[i % len(colors)]
        _add_chevron(slide, x, y_chevron, chevron_w, chevron_h, color,
                     text=step, font_size=Pt(12), font_color=D.WHITE)

        # Step number circle above
        circle_size = Inches(0.5)
        circle_x = x + (chevron_w - circle_size) // 2
        _add_oval(slide, circle_x, Inches(1.85), circle_size, circle_size,
                  color, text=str(i + 1), font_size=Pt(14), font_color=D.WHITE)

    # Description area below
    desc_y = Inches(4.2)
    desc_w = total_width / n
    for i, step in enumerate(steps):
        x = D.MARGIN_LEFT + i * desc_w
        _add_textbox(
            slide, x, desc_y, desc_w, Inches(2.5),
            text=step,
            font_size=Pt(13), font_color=D.GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_timeline_slide(prs, title, milestones, notes=""):
    """Slide 10 — Timeline: horizontal line with milestones above/below.

    milestones: list of (date_label, description) tuples
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    # Horizontal timeline line
    line_y = Inches(4.0)
    line_left = D.MARGIN_LEFT + Inches(0.3)
    line_w = D.CONTENT_WIDTH - Inches(0.6)
    _add_line(slide, line_left, line_y, line_w, Pt(4), D.NAVY)

    n = len(milestones)
    spacing = line_w / max(n - 1, 1) if n > 1 else line_w
    dot_size = Inches(0.3)

    for i, (date_label, description) in enumerate(milestones):
        if n == 1:
            x_center = line_left + line_w // 2
        else:
            x_center = line_left + int(i * spacing)

        # Dot on the line
        _add_oval(
            slide,
            x_center - dot_size // 2, line_y - dot_size // 2,
            dot_size, dot_size,
            D.ORANGE,
        )

        text_w = Inches(2.2)
        text_x = x_center - text_w // 2

        # Alternate above/below
        if i % 2 == 0:
            # Date above
            _add_textbox(
                slide, text_x, Inches(2.2), text_w, Inches(0.5),
                text=date_label,
                font_size=Pt(14), font_color=D.ORANGE,
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # Description above
            _add_textbox(
                slide, text_x, Inches(2.7), text_w, Inches(1.0),
                text=description,
                font_size=Pt(12), font_color=D.GRAY,
                alignment=PP_ALIGN.CENTER,
            )
            # Vertical connector
            _add_line(slide, x_center, Inches(3.7), Pt(2), Inches(0.3), D.LIGHT_GRAY)
        else:
            # Vertical connector
            _add_line(slide, x_center, line_y + dot_size // 2, Pt(2), Inches(0.3), D.LIGHT_GRAY)
            # Date below
            _add_textbox(
                slide, text_x, Inches(4.6), text_w, Inches(0.5),
                text=date_label,
                font_size=Pt(14), font_color=D.ORANGE,
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # Description below
            _add_textbox(
                slide, text_x, Inches(5.1), text_w, Inches(1.0),
                text=description,
                font_size=Pt(12), font_color=D.GRAY,
                alignment=PP_ALIGN.CENTER,
            )

    _add_speaker_notes(slide, notes)
    return slide


def add_matrix_slide(prs, title, top_left, top_right, bottom_left, bottom_right,
                     x_label="", y_label="", notes=""):
    """Slide 11 — 2x2 Matrix: four colored quadrants with labels.

    Each quadrant is a dict: {"title": "...", "items": ["...", "..."]}
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    # Matrix dimensions
    matrix_left = Inches(1.8)
    matrix_top = Inches(2.0)
    cell_w = Inches(4.8)
    cell_h = Inches(2.5)
    gap = Inches(0.1)
    colors = D.MATRIX_COLORS

    quadrants = [
        (0, 0, top_left, colors[0]),
        (1, 0, top_right, colors[1]),
        (0, 1, bottom_left, colors[2]),
        (1, 1, bottom_right, colors[3]),
    ]

    for col, row, data, bg_color in quadrants:
        x = matrix_left + col * (cell_w + gap)
        y = matrix_top + row * (cell_h + gap)

        # Background rectangle
        _add_rounded_rectangle(slide, x, y, cell_w, cell_h, bg_color)

        # Quadrant title
        _add_textbox(
            slide, x + Inches(0.2), y + Inches(0.15),
            cell_w - Inches(0.4), Inches(0.5),
            text=data["title"],
            font_size=Pt(16), font_color=D.NAVY,
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # Quadrant items
        if data.get("items"):
            _add_multiline_textbox(
                slide, x + Inches(0.3), y + Inches(0.7),
                cell_w - Inches(0.5), cell_h - Inches(0.9),
                lines=data["items"],
                font_size=Pt(13), font_color=D.DARK_TEXT,
                bullet_char=D.BULLET_CHAR, bullet_color=D.ORANGE,
                line_spacing=Pt(4),
            )

    # Axis labels
    if y_label:
        _add_textbox(
            slide, Inches(0.2), matrix_top + cell_h - Inches(0.3),
            Inches(1.4), Inches(0.5),
            text=y_label,
            font_size=Pt(13), font_color=D.NAVY,
            bold=True, alignment=PP_ALIGN.CENTER,
        )
    if x_label:
        _add_textbox(
            slide, matrix_left + cell_w - Inches(0.5),
            matrix_top + 2 * cell_h + gap + Inches(0.15),
            cell_w + gap, Inches(0.4),
            text=x_label,
            font_size=Pt(13), font_color=D.NAVY,
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_pyramid_slide(prs, title, levels, notes=""):
    """Slide 12 — Pyramid: stacked horizontal bars narrowing upward.

    levels: list of strings from top (smallest) to bottom (widest)
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    n = len(levels)
    pyramid_top = Inches(2.0)
    total_height = Inches(5.0)
    level_h = total_height / n
    max_width = Inches(10.0)
    min_width = Inches(3.0)
    center_x = D.SLIDE_WIDTH // 2
    colors = D.PYRAMID_COLORS

    for i, level_text in enumerate(levels):
        # Width narrows toward the top
        ratio = (n - i) / n
        w = min_width + (max_width - min_width) * ratio
        x = center_x - w // 2
        y = pyramid_top + i * level_h
        color = colors[i % len(colors)]

        _add_rounded_rectangle(
            slide, x, y, w, level_h - Inches(0.08), color,
            text=level_text, font_size=Pt(16), font_color=D.WHITE, bold=True,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_bar_chart_slide(prs, title, categories, values, notes=""):
    """Slide 13 — Bar chart: vertical bars with categories."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    # Chart
    _add_chart_bar(
        slide,
        left=D.MARGIN_LEFT + Inches(0.5), top=Inches(2.0),
        width=D.CONTENT_WIDTH - Inches(1.0), height=Inches(4.8),
        categories=categories, values=values,
    )

    _add_speaker_notes(slide, notes)
    return slide


def add_pie_chart_slide(prs, title, categories, values, notes=""):
    """Slide 14 — Pie chart: colored segments with percentages."""
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    # Pie chart
    _add_chart_pie(
        slide,
        left=Inches(2.5), top=Inches(1.8),
        width=Inches(8.0), height=Inches(5.2),
        categories=categories, values=values,
    )

    _add_speaker_notes(slide, notes)
    return slide


def add_icon_cards_slide(prs, title, cards, notes=""):
    """Slide 15 — Icon cards: grid of KPI/metric cards.

    cards: list of dicts {"value": "78%", "label": "Satisfaction", "color": RGBColor (optional)}
    Max 6 cards (2 rows x 3 cols or 1 row x 3-4).
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    n = len(cards)
    if n <= 3:
        cols, rows = n, 1
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 4, 2

    gap = Inches(0.3)
    card_w = (D.CONTENT_WIDTH - gap * (cols - 1)) / cols
    card_h = Inches(2.2) if rows == 1 else Inches(2.0)
    start_y = Inches(2.2) if rows == 1 else Inches(2.0)
    colors = D.PROCESS_COLORS

    for i, card in enumerate(cards):
        col = i % cols
        row = i // cols
        x = D.MARGIN_LEFT + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        color = card.get("color", colors[i % len(colors)])

        # Card background
        _add_rounded_rectangle(slide, x, y, card_w, card_h, D.CARD_BG,
                               border_color=D.LIGHT_GRAY)

        # Color accent bar at top of card
        _add_rectangle(slide, x, y, card_w, Inches(0.08), color)

        # Big value
        _add_textbox(
            slide, x, y + Inches(0.2),
            card_w, Inches(1.0),
            text=card["value"],
            font_size=D.CARD_TITLE_SIZE, font_color=color,
            bold=True, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.BOTTOM,
        )

        # Label
        _add_textbox(
            slide, x + Inches(0.1), y + Inches(1.3),
            card_w - Inches(0.2), Inches(0.7),
            text=card["label"],
            font_size=Pt(13), font_color=D.GRAY,
            bold=False, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.TOP,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_org_chart_slide(prs, title, manager, reports, notes=""):
    """Slide 16 — Org chart: manager node on top, direct reports below.

    manager: dict {"name": "...", "title": "..."}
    reports: list of dicts {"name": "...", "title": "..."}
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    center_x = D.SLIDE_WIDTH // 2

    # Manager node — navy rounded rectangle, centered
    mgr_w = Inches(3.0)
    mgr_h = Inches(1.0)
    mgr_x = center_x - mgr_w // 2
    mgr_y = Inches(2.0)

    _add_rounded_rectangle(
        slide, mgr_x, mgr_y, mgr_w, mgr_h, D.NAVY,
        text=f"{manager['name']}\n{manager['title']}",
        font_size=Pt(14), font_color=D.WHITE, bold=True,
    )

    # Connector: vertical line from manager bottom to horizontal bar
    n = len(reports)
    connector_y_start = mgr_y + mgr_h
    connector_y_end = Inches(3.5)
    _add_line(slide, center_x, connector_y_start, Pt(2),
              connector_y_end - connector_y_start, D.LIGHT_GRAY)

    # Report cards
    report_w = min(Inches(2.2), (D.CONTENT_WIDTH - Inches(0.2) * (n - 1)) / n) if n > 0 else Inches(2.2)
    gap = Inches(0.2)
    total_w = n * report_w + (n - 1) * gap if n > 0 else 0
    start_x = center_x - total_w // 2
    report_y = Inches(4.2)
    report_h = Inches(1.0)

    if n > 1:
        # Horizontal connector bar
        bar_left = start_x + report_w // 2
        bar_right = start_x + (n - 1) * (report_w + gap) + report_w // 2
        _add_line(slide, bar_left, connector_y_end, bar_right - bar_left, Pt(2), D.LIGHT_GRAY)

    for i, report in enumerate(reports):
        rx = start_x + i * (report_w + gap)
        rx_center = rx + report_w // 2

        # Vertical connector from bar to report card
        _add_line(slide, rx_center, connector_y_end, Pt(2),
                  report_y - connector_y_end, D.LIGHT_GRAY)

        # Report card — light gray with navy border
        _add_rounded_rectangle(
            slide, rx, report_y, report_w, report_h, D.LIGHT_GRAY,
            border_color=D.NAVY,
            text=f"{report['name']}\n{report['title']}",
            font_size=Pt(12), font_color=D.NAVY, bold=False,
        )

    _add_speaker_notes(slide, notes)
    return slide


def add_funnel_slide(prs, title, stages, notes=""):
    """Slide 17 — Funnel: centered horizontal bars decreasing in width.

    stages: list of dicts {"label": "Applied", "value": "150"}
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    n = len(stages)
    center_x = D.SLIDE_WIDTH // 2
    max_width = Inches(8.0)
    min_width = Inches(3.0)
    funnel_top = Inches(2.0)
    total_height = Inches(4.5)
    bar_h = total_height / n
    colors = D.PROCESS_COLORS

    for i, stage in enumerate(stages):
        # Width decreases linearly
        ratio = (n - i) / n if n > 1 else 1.0
        w = min_width + (max_width - min_width) * ratio
        x = center_x - w // 2
        y = funnel_top + i * bar_h
        color = colors[i % len(colors)]

        _add_rounded_rectangle(
            slide, x, y, w, bar_h - Inches(0.08), color,
        )

        # Label on the left side of the bar
        _add_textbox(
            slide, x + Inches(0.3), y,
            w // 2 - Inches(0.3), bar_h - Inches(0.08),
            text=stage["label"],
            font_size=Pt(15), font_color=D.WHITE,
            bold=True, alignment=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value on the right side of the bar
        _add_textbox(
            slide, x + w // 2, y,
            w // 2 - Inches(0.3), bar_h - Inches(0.08),
            text=stage["value"],
            font_size=Pt(15), font_color=D.WHITE,
            bold=True, alignment=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Triangle pointer at bottom
    tri_w = Inches(0.6)
    tri_h = Inches(0.4)
    tri_x = center_x - tri_w // 2
    tri_y = funnel_top + total_height + Inches(0.1)
    _add_triangle(slide, tri_x, tri_y, tri_w, tri_h, D.ORANGE)

    _add_speaker_notes(slide, notes)
    return slide


def add_team_grid_slide(prs, title, members, notes=""):
    """Slide 18 — Team grid: profile cards with initials in a grid.

    members: list of dicts {"name": "...", "role": "...", "desc": "..." (optional)}
    Max 6 members (2 rows x 3 cols).
    """
    slide = _add_blank_slide(prs)
    _set_slide_background(slide, D.WHITE)

    # Title
    _add_textbox(
        slide,
        left=D.MARGIN_LEFT, top=D.MARGIN_TOP,
        width=D.CONTENT_WIDTH, height=Inches(0.8),
        text=title,
        font_size=D.TITLE_SIZE, font_color=D.NAVY,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_line(slide, D.MARGIN_LEFT, Inches(1.5), Inches(2), Pt(3), D.ORANGE)

    n = min(len(members), 6)
    cols = min(n, 3)
    rows = 1 if n <= 3 else 2

    gap = Inches(0.4)
    card_w = (D.CONTENT_WIDTH - gap * (cols - 1)) / cols
    card_h = Inches(2.3) if rows == 1 else Inches(2.1)
    start_y = Inches(2.2) if rows == 1 else Inches(2.0)

    for i in range(n):
        member = members[i]
        col = i % cols
        row = i // cols
        x = D.MARGIN_LEFT + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        # Card background
        _add_rounded_rectangle(slide, x, y, card_w, card_h, D.LIGHT_GRAY,
                               border_color=D.LIGHT_GRAY)

        # Initials circle
        circle_size = Inches(0.7)
        circle_x = x + (card_w - circle_size) // 2
        circle_y = y + Inches(0.2)

        # Extract initials from name
        parts = member["name"].split()
        initials = "".join(p[0].upper() for p in parts[:2]) if parts else "?"

        _add_oval(
            slide, circle_x, circle_y, circle_size, circle_size,
            D.WHITE, text=initials, font_size=Pt(18), font_color=D.NAVY, bold=True,
        )

        # Name
        _add_textbox(
            slide, x, circle_y + circle_size + Inches(0.1),
            card_w, Inches(0.4),
            text=member["name"],
            font_size=Pt(14), font_color=D.NAVY,
            bold=True, alignment=PP_ALIGN.CENTER,
        )

        # Role
        _add_textbox(
            slide, x, circle_y + circle_size + Inches(0.45),
            card_w, Inches(0.35),
            text=member["role"],
            font_size=Pt(12), font_color=D.ORANGE,
            bold=True, alignment=PP_ALIGN.CENTER,
        )

        # Description (optional)
        desc = member.get("desc", "")
        if desc:
            _add_textbox(
                slide, x + Inches(0.1), circle_y + circle_size + Inches(0.8),
                card_w - Inches(0.2), Inches(0.5),
                text=desc,
                font_size=Pt(10), font_color=D.GRAY,
                bold=False, alignment=PP_ALIGN.CENTER,
            )

    _add_speaker_notes(slide, notes)
    return slide
