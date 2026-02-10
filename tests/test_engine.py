"""Unit tests for slide_engine — each layout individually."""

import os
import sys
import pytest

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_engine import (
    create_presentation,
    save_presentation,
    add_title_slide,
    add_agenda_slide,
    add_section_slide,
    add_bullets_slide,
    add_two_columns_slide,
    add_key_stat_slide,
    add_quote_slide,
    add_conclusion_slide,
    add_process_flow_slide,
    add_timeline_slide,
    add_matrix_slide,
    add_pyramid_slide,
    add_bar_chart_slide,
    add_pie_chart_slide,
    add_icon_cards_slide,
    add_org_chart_slide,
    add_funnel_slide,
    add_team_grid_slide,
)


@pytest.fixture
def prs():
    return create_presentation()


class TestCreatePresentation:
    def test_creates_presentation(self):
        prs = create_presentation()
        assert prs is not None
        assert len(prs.slides) == 0

    def test_slide_dimensions_16_9(self):
        from pptx.util import Inches
        prs = create_presentation()
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)


class TestSavePresentation:
    def test_save_adds_extension(self, prs, tmp_path):
        filepath = str(tmp_path / "test")
        result = save_presentation(prs, filepath)
        assert result.endswith(".pptx")
        assert os.path.exists(result)

    def test_save_keeps_extension(self, prs, tmp_path):
        filepath = str(tmp_path / "test.pptx")
        result = save_presentation(prs, filepath)
        assert result == filepath
        assert os.path.exists(result)


class TestTitleSlide:
    def test_basic(self, prs):
        slide = add_title_slide(prs, "Mon Titre", "Sous-titre")
        assert len(prs.slides) == 1
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert "Mon Titre" in texts

    def test_without_subtitle(self, prs):
        slide = add_title_slide(prs, "Titre seul")
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_title_slide(prs, "Titre", notes="Notes du présentateur")
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Notes du présentateur" in notes_text

    def test_unicode(self, prs):
        slide = add_title_slide(prs, "Données RH : évaluation & développement")
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert any("évaluation" in t for t in texts)


class TestAgendaSlide:
    def test_basic(self, prs):
        items = ["Introduction", "Diagnostic", "Préconisations"]
        slide = add_agenda_slide(prs, items)
        assert len(prs.slides) == 1

    def test_custom_title(self, prs):
        slide = add_agenda_slide(prs, ["A", "B"], title="Plan")
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert "Plan" in texts

    def test_with_notes(self, prs):
        slide = add_agenda_slide(prs, ["A"], notes="Rappeler le contexte")
        assert "Rappeler le contexte" in slide.notes_slide.notes_text_frame.text


class TestSectionSlide:
    def test_basic(self, prs):
        slide = add_section_slide(prs, "Partie I")
        assert len(prs.slides) == 1

    def test_with_subtitle(self, prs):
        slide = add_section_slide(prs, "Partie I", "Cadre théorique")
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert "Cadre théorique" in texts


class TestBulletsSlide:
    def test_basic(self, prs):
        slide = add_bullets_slide(prs, "Points clés", ["Point A", "Point B", "Point C"])
        assert len(prs.slides) == 1

    def test_unicode_bullets(self, prs):
        slide = add_bullets_slide(prs, "Résumé", ["Première étape", "Deuxième étape"])
        assert len(prs.slides) == 1


class TestTwoColumnsSlide:
    def test_basic(self, prs):
        slide = add_two_columns_slide(
            prs,
            title="Comparaison",
            left_title="Avant",
            left_items=["Item A", "Item B"],
            right_title="Après",
            right_items=["Item C", "Item D"],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_two_columns_slide(
            prs,
            title="Test",
            left_title="L", left_items=["a"],
            right_title="R", right_items=["b"],
            notes="Comparer les deux approches",
        )
        assert "Comparer" in slide.notes_slide.notes_text_frame.text


class TestKeyStatSlide:
    def test_basic(self, prs):
        slide = add_key_stat_slide(prs, "78%", "des salariés satisfaits")
        assert len(prs.slides) == 1

    def test_stat_text(self, prs):
        slide = add_key_stat_slide(prs, "3.2M€", "Budget formation annuel")
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert "3.2M€" in texts


class TestQuoteSlide:
    def test_basic(self, prs):
        slide = add_quote_slide(prs, "La GRH est un levier stratégique.")
        assert len(prs.slides) == 1

    def test_with_author(self, prs):
        slide = add_quote_slide(
            prs,
            "Le capital humain est la première richesse.",
            author="Jean-Marie Peretti",
        )
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert any("Peretti" in t for t in texts)


class TestConclusionSlide:
    def test_basic(self, prs):
        slide = add_conclusion_slide(
            prs,
            "Conclusion",
            ["Mettre en place un SIRH", "Former les managers", "Suivre les KPIs"],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_conclusion_slide(
            prs, "À retenir", ["Point 1"], notes="Ouvrir sur les perspectives"
        )
        assert "perspectives" in slide.notes_slide.notes_text_frame.text


class TestSpeakerNotes:
    def test_empty_notes_no_crash(self, prs):
        slide = add_title_slide(prs, "Test", notes="")
        # Should not crash with empty notes

    def test_long_notes(self, prs):
        long_notes = "Note très longue. " * 100
        slide = add_title_slide(prs, "Test", notes=long_notes)
        assert len(slide.notes_slide.notes_text_frame.text) > 100

    def test_unicode_notes(self, prs):
        slide = add_bullets_slide(
            prs, "Test", ["A"],
            notes="Référence : Thévenet (2015), « Fonctions RH » — 4ème édition"
        )
        assert "Thévenet" in slide.notes_slide.notes_text_frame.text


# ===================================================================
# VISUAL LAYOUTS TESTS
# ===================================================================


class TestProcessFlowSlide:
    def test_basic(self, prs):
        slide = add_process_flow_slide(
            prs, "Processus GPEC",
            ["Diagnostic", "Planification", "Action", "Évaluation"],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_process_flow_slide(
            prs, "Processus", ["A", "B", "C"],
            notes="Décrire chaque étape",
        )
        assert "Décrire" in slide.notes_slide.notes_text_frame.text

    def test_many_steps(self, prs):
        slide = add_process_flow_slide(
            prs, "Long process",
            ["Étape 1", "Étape 2", "Étape 3", "Étape 4", "Étape 5", "Étape 6"],
        )
        assert len(prs.slides) == 1


class TestTimelineSlide:
    def test_basic(self, prs):
        slide = add_timeline_slide(
            prs, "Évolution GPEC",
            [("2005", "Loi Borloo"), ("2013", "ANI"), ("2017", "Ordonnances Macron"), ("2024", "GEPP 2.0")],
        )
        assert len(prs.slides) == 1

    def test_single_milestone(self, prs):
        slide = add_timeline_slide(
            prs, "Date clé",
            [("2025", "Mise en place SIRH")],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_timeline_slide(
            prs, "Timeline",
            [("T1", "Action A"), ("T2", "Action B")],
            notes="Respecter le phasage",
        )
        assert "phasage" in slide.notes_slide.notes_text_frame.text


class TestMatrixSlide:
    def test_basic_swot(self, prs):
        slide = add_matrix_slide(
            prs, "Analyse SWOT",
            top_left={"title": "Forces", "items": ["Culture forte", "Budget"]},
            top_right={"title": "Faiblesses", "items": ["Pas de SIRH"]},
            bottom_left={"title": "Opportunités", "items": ["IA", "GEPP"]},
            bottom_right={"title": "Menaces", "items": ["Turnover"]},
        )
        assert len(prs.slides) == 1

    def test_with_axis_labels(self, prs):
        slide = add_matrix_slide(
            prs, "Matrice compétences",
            top_left={"title": "Q1", "items": ["A"]},
            top_right={"title": "Q2", "items": ["B"]},
            bottom_left={"title": "Q3", "items": ["C"]},
            bottom_right={"title": "Q4", "items": ["D"]},
            x_label="Criticité →",
            y_label="Maîtrise →",
        )
        assert len(prs.slides) == 1


class TestPyramidSlide:
    def test_basic(self, prs):
        slide = add_pyramid_slide(
            prs, "Pyramide de Maslow",
            ["Accomplissement", "Estime", "Appartenance", "Sécurité", "Physiologie"],
        )
        assert len(prs.slides) == 1

    def test_three_levels(self, prs):
        slide = add_pyramid_slide(
            prs, "Niveaux de compétences",
            ["Expert", "Confirmé", "Junior"],
        )
        assert len(prs.slides) == 1


class TestBarChartSlide:
    def test_basic(self, prs):
        slide = add_bar_chart_slide(
            prs, "Taux de turnover par département",
            categories=["RH", "IT", "Finance", "Marketing", "Production"],
            values=[8.5, 15.2, 6.3, 12.1, 9.8],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_bar_chart_slide(
            prs, "KPI Formation",
            categories=["2022", "2023", "2024"],
            values=[45, 62, 78],
            notes="Progression constante du taux de formation",
        )
        assert "Progression" in slide.notes_slide.notes_text_frame.text


class TestPieChartSlide:
    def test_basic(self, prs):
        slide = add_pie_chart_slide(
            prs, "Répartition des effectifs",
            categories=["CDI", "CDD", "Intérim", "Alternance"],
            values=[65, 15, 12, 8],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_pie_chart_slide(
            prs, "Budget formation",
            categories=["Technique", "Management", "Soft skills"],
            values=[50, 30, 20],
            notes="Rééquilibrer vers les soft skills",
        )
        assert "soft skills" in slide.notes_slide.notes_text_frame.text


class TestIconCardsSlide:
    def test_basic_three_cards(self, prs):
        slide = add_icon_cards_slide(
            prs, "Dashboard RH",
            cards=[
                {"value": "78%", "label": "Satisfaction"},
                {"value": "12%", "label": "Turnover"},
                {"value": "3.2j", "label": "Absentéisme"},
            ],
        )
        assert len(prs.slides) == 1

    def test_six_cards(self, prs):
        slide = add_icon_cards_slide(
            prs, "KPIs RH",
            cards=[
                {"value": "156", "label": "Effectif total"},
                {"value": "42%", "label": "Femmes managers"},
                {"value": "2.1%", "label": "Budget formation / MS"},
                {"value": "8.5%", "label": "Turnover"},
                {"value": "92%", "label": "Entretiens réalisés"},
                {"value": "4.2/5", "label": "Score engagement"},
            ],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_icon_cards_slide(
            prs, "Indicateurs",
            cards=[{"value": "95%", "label": "Taux complétion"}],
            notes="Objectif dépassé",
        )
        assert "dépassé" in slide.notes_slide.notes_text_frame.text


class TestOrgChartSlide:
    def test_basic(self, prs):
        slide = add_org_chart_slide(
            prs, "Organisation RH",
            manager={"name": "Marie Dupont", "title": "DRH"},
            reports=[
                {"name": "Jean Martin", "title": "Resp. Formation"},
                {"name": "Sophie Leclerc", "title": "Resp. Recrutement"},
                {"name": "Pierre Durand", "title": "Resp. Paie"},
            ],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_org_chart_slide(
            prs, "Organigramme",
            manager={"name": "A", "title": "DRH"},
            reports=[{"name": "B", "title": "Adj"}],
            notes="Présenter l'équipe",
        )
        assert "équipe" in slide.notes_slide.notes_text_frame.text

    def test_single_report(self, prs):
        slide = add_org_chart_slide(
            prs, "Hiérarchie",
            manager={"name": "Chef", "title": "Directeur"},
            reports=[{"name": "Adjoint", "title": "Sous-directeur"}],
        )
        assert len(prs.slides) == 1


class TestFunnelSlide:
    def test_basic(self, prs):
        slide = add_funnel_slide(
            prs, "Tunnel de recrutement",
            stages=[
                {"label": "Candidatures", "value": "150"},
                {"label": "Entretiens RH", "value": "45"},
                {"label": "Entretiens manager", "value": "20"},
                {"label": "Offres", "value": "8"},
                {"label": "Embauches", "value": "5"},
            ],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_funnel_slide(
            prs, "Funnel",
            stages=[{"label": "A", "value": "100"}, {"label": "B", "value": "50"}],
            notes="Analyser le taux de conversion",
        )
        assert "conversion" in slide.notes_slide.notes_text_frame.text

    def test_two_stages(self, prs):
        slide = add_funnel_slide(
            prs, "Mini funnel",
            stages=[
                {"label": "Candidatures", "value": "200"},
                {"label": "Recrutés", "value": "10"},
            ],
        )
        assert len(prs.slides) == 1


class TestTeamGridSlide:
    def test_basic(self, prs):
        slide = add_team_grid_slide(
            prs, "Équipe projet",
            members=[
                {"name": "Marie Dupont", "role": "Chef de projet", "desc": "Pilotage GPEC"},
                {"name": "Jean Martin", "role": "RRH", "desc": "Référent compétences"},
                {"name": "Sophie Leclerc", "role": "Consultante"},
            ],
        )
        assert len(prs.slides) == 1

    def test_six_members(self, prs):
        slide = add_team_grid_slide(
            prs, "Équipe complète",
            members=[
                {"name": "A B", "role": "R1", "desc": "D1"},
                {"name": "C D", "role": "R2", "desc": "D2"},
                {"name": "E F", "role": "R3", "desc": "D3"},
                {"name": "G H", "role": "R4", "desc": "D4"},
                {"name": "I J", "role": "R5", "desc": "D5"},
                {"name": "K L", "role": "R6", "desc": "D6"},
            ],
        )
        assert len(prs.slides) == 1

    def test_with_notes(self, prs):
        slide = add_team_grid_slide(
            prs, "Équipe",
            members=[{"name": "Test User", "role": "Dev"}],
            notes="Présenter les rôles",
        )
        assert "rôles" in slide.notes_slide.notes_text_frame.text

    def test_without_desc(self, prs):
        slide = add_team_grid_slide(
            prs, "Équipe sans description",
            members=[
                {"name": "Alice Bonnet", "role": "Manager"},
                {"name": "Bob Petit", "role": "Analyste"},
            ],
        )
        assert len(prs.slides) == 1
